import sys
import os
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

target_dir = "/Users/howardliao/Desktop/Howard/Howard_CISO"
photo_path = "/Users/howardliao/Desktop/Howard/Howard_CISO/assets/howard_portrait.jpg"
asset_dir = "/Users/howardliao/Desktop/Howard/Howard_CISO/assets"

# Initialize 16:9 Presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color Palette
BG_DARK = RGBColor(2, 6, 23)        # #020617 Deep Obsidian
BG_CARD = RGBColor(15, 23, 42)      # #0F172A Card Background
BG_CARD_LIGHT = RGBColor(30, 41, 59)# #1E293B Lighter Container
BORDER_CYAN = RGBColor(14, 165, 233)# #0EA5E9
TEXT_WHITE = RGBColor(255, 255, 255)
TEXT_LIGHT = RGBColor(226, 232, 240)# #E2E8F0
TEXT_MUTED = RGBColor(148, 163, 184)# #94A3B8
COLOR_CYAN = RGBColor(56, 189, 248) # #38BDF8
COLOR_GOLD = RGBColor(245, 158, 11) # #F59E0B
COLOR_EMERALD = RGBColor(16, 185, 129) # #10B981
COLOR_PURPLE = RGBColor(168, 85, 247)  # #A855F7
COLOR_BLUE = RGBColor(59, 130, 246)    # #3B82F6

FONT_TITLE = "Microsoft JhengHei"
FONT_BODY = "Microsoft JhengHei"

def set_slide_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK

def add_header(slide, title_text, category_tag="STRATEGIC EXECUTIVE BRIEFING"):
    tb_tag = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.3))
    p_tag = tb_tag.text_frame.paragraphs[0]
    p_tag.text = f"🛡️ {category_tag}"
    p_tag.font.name = FONT_TITLE
    p_tag.font.size = Pt(11)
    p_tag.font.bold = True
    p_tag.font.color.rgb = COLOR_CYAN

    tb_title = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.5), Inches(0.6))
    p_title = tb_title.text_frame.paragraphs[0]
    p_title.text = title_text
    p_title.font.name = FONT_TITLE
    p_title.font.size = Pt(22)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_WHITE

    shape_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.8), Inches(1.35), Inches(11.7), Inches(0.02)
    )
    shape_line.fill.solid()
    shape_line.fill.fore_color.rgb = BORDER_CYAN
    shape_line.line.color.rgb = BORDER_CYAN

# ==========================================
# SLIDE 1: COVER SLIDE
# ==========================================
slide_layout = prs.slide_layouts[6]
slide1 = prs.slides.add_slide(slide_layout)
set_slide_background(slide1)

accent_box = slide1.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.8), Inches(0.8), Inches(7.5), Inches(5.9)
)
accent_box.fill.solid()
accent_box.fill.fore_color.rgb = BG_CARD
accent_box.line.color.rgb = BORDER_CYAN
accent_box.line.width = Pt(1.5)

tb_aud = slide1.shapes.add_textbox(Inches(1.2), Inches(1.1), Inches(6.8), Inches(0.4))
p_aud = tb_aud.text_frame.paragraphs[0]
p_aud.text = "🎯 面試專用簡報（建議時間：10–12 分鐘）| 呈報對象：董事長、總經理／CEO、HR Director 及高階決策層"
p_aud.font.name = FONT_TITLE
p_aud.font.size = Pt(10)
p_aud.font.bold = True
p_aud.font.color.rgb = COLOR_CYAN

tb_c_title = slide1.shapes.add_textbox(Inches(1.2), Inches(1.6), Inches(6.8), Inches(1.8))
p_c1 = tb_c_title.text_frame.paragraphs[0]
p_c1.text = "從資安韌性到商業賦能"
p_c1.font.name = FONT_TITLE
p_c1.font.size = Pt(28)
p_c1.font.bold = True
p_c1.font.color.rgb = TEXT_WHITE

p_c2 = tb_c_title.text_frame.add_paragraph()
p_c2.text = "跨國集團全球科技治理與高可用架構實戰"
p_c2.font.name = FONT_TITLE
p_c2.font.size = Pt(20)
p_c2.font.bold = True
p_c2.font.color.rgb = COLOR_GOLD

tb_c_sub = slide1.shapes.add_textbox(Inches(1.2), Inches(3.5), Inches(6.8), Inches(2.8))
tf_c = tb_c_sub.text_frame
tf_c.word_wrap = True

bullets_cover = [
    "🏛️ 董事會級 CISO 戰略治理：將資安風險轉化為具體 ROI、TCO 與營運韌性指標",
    "☁️ 多雲原生高可用架構：GKE + MongoDB Atlas 實現 100% 零停機 (Zero Outage)",
    "💰 多雲財務營運 (FinOps)：兼顧極致效能與穩定性，實現雲端成本壓降 30%",
    "🤖 前瞻 AI 治理體系：基於 ISO/IEC 42001 建立 GenAI 資料防洩與 AISDLC 護城河"
]
for b in bullets_cover:
    p = tf_c.add_paragraph()
    p.text = b
    p.font.name = FONT_BODY
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_LIGHT
    p.space_after = Pt(8)

right_card = slide1.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(8.6), Inches(0.8), Inches(3.9), Inches(5.9)
)
right_card.fill.solid()
right_card.fill.fore_color.rgb = BG_CARD_LIGHT
right_card.line.color.rgb = BORDER_CYAN
right_card.line.width = Pt(1)

if os.path.exists(photo_path):
    slide1.shapes.add_picture(photo_path, Inches(9.45), Inches(1.1), width=Inches(2.2))

tb_spk = slide1.shapes.add_textbox(Inches(8.8), Inches(4.3), Inches(3.5), Inches(2.2))
tf_spk = tb_spk.text_frame
tf_spk.word_wrap = True

p_s1 = tf_spk.paragraphs[0]
p_s1.text = "廖倫豪 博士 (Howard Liao, Ph.D.)"
p_s1.font.name = FONT_TITLE
p_s1.font.size = Pt(16)
p_s1.font.bold = True
p_s1.font.color.rgb = TEXT_WHITE
p_s1.alignment = PP_ALIGN.CENTER

p_s2 = tf_spk.add_paragraph()
p_s2.text = "集團資安長 暨 科技副總 (Group CISO)"
p_s2.font.name = FONT_BODY
p_s2.font.size = Pt(11)
p_s2.font.bold = True
p_s2.font.color.rgb = COLOR_CYAN
p_s2.alignment = PP_ALIGN.CENTER

p_s3 = tf_spk.add_paragraph()
p_s3.text = "朝陽科大資管博士 | 27年IT領導 | 15年資安戰略\nISO 27001 & ISO 42001 主任稽核員 | PMP"
p_s3.font.name = FONT_BODY
p_s3.font.size = Pt(9.5)
p_s3.font.color.rgb = TEXT_MUTED
p_s3.alignment = PP_ALIGN.CENTER
p_s3.space_before = Pt(4)

# ==========================================
# SLIDE 2: 學經歷與領導底蘊
# ==========================================
slide2 = prs.slides.add_slide(slide_layout)
set_slide_background(slide2)
add_header(slide2, "一、學經歷簡介：27 年 IT 管理與 15 年資安戰略底蘊", "EXECUTIVE TRAJECTORY & CREDENTIALS")

col1 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(3.7), Inches(5.3))
col1.fill.solid(); col1.fill.fore_color.rgb = BG_CARD; col1.line.color.rgb = COLOR_PURPLE

tb1 = slide2.shapes.add_textbox(Inches(0.9), Inches(1.7), Inches(3.5), Inches(5.1))
tf1 = tb1.text_frame; tf1.word_wrap = True
p = tf1.paragraphs[0]; p.text = "🎓 學術背景與國際專業認證"; p.font.name = FONT_TITLE; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = COLOR_PURPLE

items1 = [
    "🏛️ 朝陽科技大學 資訊科技管理博士 (Ph.D.)\n   研究：SOA 服務導向架構與 IT 治理模型",
    "📚 國立中興大學 資訊科學碩士 (M.S.)\n   專注：現代密碼學、數據治理、軟體工程",
    "💻 國立中興大學 資訊科學學士 (B.S.)",
    "🛡️ ISO 27001 資安主任稽核員 (Lead Auditor)",
    "🤖 ISO/IEC 42001 (AIMS) AI 治理主任稽核員",
    "📊 國際專案管理師 (PMP) & 敏捷大師 (CSM)",
    "🗄️ Oracle 認證專家 (OCP) & ESG 永續規劃師"
]
for item in items1:
    p = tf1.add_paragraph()
    p.text = item
    p.font.name = FONT_BODY; p.font.size = Pt(10.5); p.font.color.rgb = TEXT_LIGHT
    p.space_after = Pt(6)

col2 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.7), Inches(1.6), Inches(4.3), Inches(5.3))
col2.fill.solid(); col2.fill.fore_color.rgb = BG_CARD; col2.line.color.rgb = COLOR_CYAN

tb2 = slide2.shapes.add_textbox(Inches(4.8), Inches(1.7), Inches(4.1), Inches(5.1))
tf2 = tb2.text_frame; tf2.word_wrap = True
p = tf2.paragraphs[0]; p.text = "🏢 跨國與上市櫃職涯里程碑 (27+ Y)"; p.font.name = FONT_TITLE; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = COLOR_CYAN

items2 = [
    "👑 2025–Present: 集團副總 暨 資安與轉型負責人\n   主導跨國集團資安策略、AI 治理，管 $12M 預算",
    "🎮 2022–2025: 隆中網絡 GameSparcs IT Director\n   百萬級手遊多雲架構，創 100% 零停機，降本 30%",
    "🏭 2018–2022: 泓晏科技 IT Director (深港台跨區)\n   電子製造 IT/OT 融合、ERP/MES/PLM 安全整合",
    "🌐 2014–2018: 凌網科技 (上市櫃) 資訊部經理\n   亞太 5 大據點 ITSM、資料中心高可用戰略",
    "🏬 2011–2014: 光南集團 (上市櫃) 資訊經理/轉型主管\n   跨國 M&A 數據清洗、ERP/POS/CRM 中台整併",
    "💻 2008–2011: 美商寶藍 (Borland) 資訊部經理",
    "🗄️ 2002–2008: 美商賽貝斯 (Sybase) 資訊部經理"
]
for item in items2:
    p = tf2.add_paragraph()
    p.text = item
    p.font.name = FONT_BODY; p.font.size = Pt(10); p.font.color.rgb = TEXT_LIGHT
    p.space_after = Pt(5)

col3 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.2), Inches(1.6), Inches(3.3), Inches(5.3))
col3.fill.solid(); col3.fill.fore_color.rgb = BG_CARD; col3.line.color.rgb = COLOR_GOLD

tb3 = slide2.shapes.add_textbox(Inches(9.3), Inches(1.7), Inches(3.1), Inches(5.1))
tf3 = tb3.text_frame; tf3.word_wrap = True
p = tf3.paragraphs[0]; p.text = "⭐ 核心優勢與量化成效"; p.font.name = FONT_TITLE; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = COLOR_GOLD

items3 = [
    "🛡️ 100% Zero Outage\n   跨區域容災，百萬並發零中斷",
    "📉 FinOps 降本 30%\n   多雲 TCO 精準治理，年省數百萬",
    "⚡ 重大事件 -30% & MTTR -30%\n   集中式 SIEM/EDR 全域可觀測",
    "🏆 雲端架構卓越獎\n   獲 Google Cloud 官方影音專訪",
    "📰 權威媒體專題實名肯定\n   CIO Taiwan / iThome 封面報導",
    "📜 SCI 頂級期刊 DOI 出版品\n   可逆秘密影像共享密碼學技術"
]
for item in items3:
    p = tf3.add_paragraph()
    p.text = item
    p.font.name = FONT_BODY; p.font.size = Pt(10); p.font.color.rgb = TEXT_LIGHT
    p.space_after = Pt(6)

# ==========================================
# SLIDE 3: 代表性案例與商業挑戰
# ==========================================
slide3 = prs.slides.add_slide(slide_layout)
set_slide_background(slide3)
add_header(slide3, "二、代表性案例：全球跨區高並發線上娛樂平台架構轉型", "CASE STUDY — BACKGROUND & BUSINESS DRIVERS")

banner = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(11.7), Inches(1.1))
banner.fill.solid(); banner.fill.fore_color.rgb = BG_CARD; banner.line.color.rgb = COLOR_CYAN

tb_b = slide3.shapes.add_textbox(Inches(1.0), Inches(1.65), Inches(11.3), Inches(1.0))
tf_b = tb_b.text_frame; tf_b.word_wrap = True
p = tf_b.paragraphs[0]; p.text = "📌 專案背景 (Project Background)"; p.font.name = FONT_TITLE; p.font.size = Pt(13); p.font.bold = True; p.font.color.rgb = COLOR_CYAN
p2 = tf_b.add_paragraph()
p2.text = "主導上市櫃跨國數位娛樂平台（隆中網絡 GameSparcs & 關聯集團）全球多雲與資安架構現代化重構。服務橫跨台、美、澳、歐、東南亞數百萬活躍玩家，需承載 24/7 全天候即時連線、高頻交易與全球跨區高並發讀寫。"
p2.font.name = FONT_BODY; p2.font.size = Pt(11); p2.font.color.rgb = TEXT_LIGHT

drivers = [
    ("🚀 商業連續性 (Business Continuity)", "• 玩家即時對弈不容許任何卡頓或中斷\n• 促銷與節慶活動流量瞬間暴增 5–10 倍\n• 舊有架構擴展彈性不足，面臨單點故障 (SPOF) 風險", COLOR_EMERALD),
    ("🛡️ 嚴峻資安威脅 (External Threat Surface)", "• 頻繁面臨 DDoS 流量攻擊與防不勝防的撞庫\n• 惡意 Bot 搶佔資源與 API 偽冒連線風險\n• 跨國營運須符合嚴格資料隱私與合規稽核", COLOR_PURPLE),
    ("💰 成本與維運失控 (Cost & Operational Load)", "• 多公有雲環境點對點架構混亂，TCO 急遽攀升\n• 缺乏集中式可觀測性，故障排查 MTTR 耗時數小時\n• 跨國維運團隊權責不清，跨區協作成本過高", COLOR_GOLD)
]

left_pos = 0.8
for title, desc, color in drivers:
    box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left_pos), Inches(2.9), Inches(3.75), Inches(3.9))
    box.fill.solid(); box.fill.fore_color.rgb = BG_CARD_LIGHT; box.line.color.rgb = color; box.line.width = Pt(1.5)
    
    tb = slide3.shapes.add_textbox(Inches(left_pos + 0.15), Inches(3.1), Inches(3.45), Inches(3.5))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = title; p.font.name = FONT_TITLE; p.font.size = Pt(13); p.font.bold = True; p.font.color.rgb = color
    
    p_d = tf.add_paragraph()
    p_d.text = desc; p_d.font.name = FONT_BODY; p_d.font.size = Pt(11); p_d.font.color.rgb = TEXT_LIGHT; p_d.space_before = Pt(8)
    left_pos += 3.98

# ==========================================
# SLIDE 4: 問題剖析與推動方法論
# ==========================================
slide4 = prs.slides.add_slide(slide_layout)
set_slide_background(slide4)
add_header(slide4, "三、問題分析與推動方法：以架構委員會與 DevSecOps 系統化落地", "SYSTEMATIC PROBLEM ANALYSIS & METHODOLOGY")

left_box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(5.7), Inches(5.3))
left_box.fill.solid(); left_box.fill.fore_color.rgb = BG_CARD; left_box.line.color.rgb = COLOR_PURPLE

tb_l = slide4.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(5.3), Inches(5.1))
tf_l = tb_l.text_frame; tf_l.word_wrap = True
p = tf_l.paragraphs[0]; p.text = "🔍 核心痛點深入剖析 (Root Cause Analysis)"; p.font.name = FONT_TITLE; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = COLOR_PURPLE

problems = [
    ("1. 架構僵化與單點失效", "過去仰賴傳統虛擬機與地端主機，無法實現跨可用區 (Multi-Zone) 與跨區域 (Multi-Region) 自動容錯切換。"),
    ("2. 資安控制點分散與影子系統", "各產品團隊各自為政，缺乏統一身分鑑別 (SSO/MFA)、特權存取 (PAM) 與 API 安全規範，攻擊面暴露過大。"),
    ("3. 數據中台併發讀寫瓶頸", "關聯式資料庫在全球百萬連線下鎖表延遲，資料庫調優已達物理極限，急需全託管分散式架構。"),
    ("4. 雲端資源無節制開支", "缺乏 FinOps 財務營運思維，過度配置 (Over-provisioning) 資源導致每月雲端帳單失控浪費。")
]
for p_t, p_d in problems:
    p = tf_l.add_paragraph()
    p.text = p_t; p.font.name = FONT_TITLE; p.font.size = Pt(11.5); p.font.bold = True; p.font.color.rgb = COLOR_CYAN; p.space_before = Pt(6)
    p_desc = tf_l.add_paragraph()
    p_desc.text = p_d; p_desc.font.name = FONT_BODY; p_desc.font.size = Pt(10.5); p_desc.font.color.rgb = TEXT_LIGHT

right_box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.3))
right_box.fill.solid(); right_box.fill.fore_color.rgb = BG_CARD; right_box.line.color.rgb = COLOR_CYAN

tb_r = slide4.shapes.add_textbox(Inches(7.0), Inches(1.7), Inches(5.3), Inches(5.1))
tf_r = tb_r.text_frame; tf_r.word_wrap = True
p = tf_r.paragraphs[0]; p.text = "⚙️ 系統化規劃與推動方法論 (Methodology)"; p.font.name = FONT_TITLE; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = COLOR_CYAN

solutions = [
    ("🏛️ 成立架構審查委員會 (ARB)", "建立 Clean Core 原則與 API-first 架構標準，所有重大技術選型與整合強制通過 ARB 審查。"),
    ("🛡️ 導入 DevSecOps Quality Gates", "在 CI/CD 流水線嵌入 SAST (靜態掃描)、DAST (動態掃描)、SCA (開源合規) 與 SBOM，漏洞左移 (Shift-Left)。"),
    ("📊 落地 OGSM 目標管理與 FinOps 矩陣", "將 100% 零停機、-30% 成本、-30% 事件轉化為跨部門 KPI，建立單位成本模型與每週用量監控。"),
    ("🔄 紅藍軍對抗與混沌工程 (Chaos Engineering)", "定期執行勒索軟體攻擊演練、跨區斷網演習與即時切換驗證，確保 RTO < 15 分鐘、RPO ≈ 0。")
]
for s_t, s_d in solutions:
    p = tf_r.add_paragraph()
    p.text = s_t; p.font.name = FONT_TITLE; p.font.size = Pt(11.5); p.font.bold = True; p.font.color.rgb = COLOR_GOLD; p.space_before = Pt(6)
    p_desc = tf_r.add_paragraph()
    p_desc.text = s_d; p_desc.font.name = FONT_BODY; p_desc.font.size = Pt(10.5); p_desc.font.color.rgb = TEXT_LIGHT

# ==========================================
# SLIDE 5: 資訊策略與 4 層整體架構設計
# ==========================================
slide5 = prs.slides.add_slide(slide_layout)
set_slide_background(slide5)
add_header(slide5, "四、資訊策略與整體架構：從業務需求延伸至 4 層防禦與中台藍圖", "ENTERPRISE ARCHITECTURE — 4-LAYER BLUEPRINT")

layers = [
    ("L1. 身分鑑別與邊界防禦層 (Zero Trust & Edge Security)", 
     "• 全面導入企業級 IAM / SSO / MFA 與條件式存取 (Conditional Access)，落實最小特權原則 (Least Privilege)\n• 部署 Cloudflare WAAP、DDoS 智慧緩解、Anti-Bot 爬蟲過濾與 API Gateway 鑑權，外部攻擊面收斂 80%",
     COLOR_CYAN),
    ("L2. 計算與微服務容器中台層 (GKE Cloud-Native Core)",
     "• 以 Google Kubernetes Engine (GKE) 為核心建構 Multi-Zone / Multi-Region 高可用叢集\n• 實現 HPA / VPA 毫秒級自動擴縮容，無縫因應 5–10 倍瞬時流量激增，確保 100% Zero Outage 零停機",
     COLOR_BLUE),
    ("L3. 數據中台與分散式資料庫層 (Data Fabric & Storage)",
     "• 引進 MongoDB Atlas 全託管跨區域多主節點叢集，支撐全球百萬玩家高並發即時對弈（可用性達 99.995%）\n• 結合 Kafka / EDI 異步事件串流中台與 Canonical Data Model，消除點對點耦合混亂",
     COLOR_EMERALD),
    ("L4. 安全維運、可觀測性與災備層 (SecOps & Cyber Resilience)",
     "• 建置集中式 SOC/SIEM、EDR/XDR、ELK Stack 與 Prometheus，即時關聯分析使 MTTR 縮短 30%\n• 實施不可變備份 (Immutable Backup)、跨公有雲冷熱備援與自動化容錯切換 (Failover)",
     COLOR_PURPLE)
]

top_pos = 1.55
for l_title, l_desc, l_color in layers:
    box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(top_pos), Inches(11.7), Inches(1.25))
    box.fill.solid(); box.fill.fore_color.rgb = BG_CARD; box.line.color.rgb = l_color; box.line.width = Pt(1.5)
    
    tb = slide5.shapes.add_textbox(Inches(1.0), Inches(top_pos + 0.08), Inches(11.3), Inches(1.1))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = l_title; p.font.name = FONT_TITLE; p.font.size = Pt(13); p.font.bold = True; p.font.color.rgb = l_color
    
    p_d = tf.add_paragraph()
    p_d.text = l_desc; p_d.font.name = FONT_BODY; p_d.font.size = Pt(10.5); p_d.font.color.rgb = TEXT_LIGHT; p_d.space_before = Pt(3)
    top_pos += 1.35

# ==========================================
# SLIDE 6: 前瞻 AI 治理與安全創新 (ISO 42001)
# ==========================================
slide6 = prs.slides.add_slide(slide_layout)
set_slide_background(slide6)
add_header(slide6, "五、前瞻創新：ISO/IEC 42001 AI 治理體系與 AISDLC 落地", "AI GOVERNANCE & SECURE INNOVATION")

ai_pillars = [
    ("🛡️ 1. ISO 42001 AI 管理體系與政策規範", 
     "• 制定全集團生成式 AI 使用辦法與審批流程\n• 明訂「五大絕對禁止輸入項目」：認證憑證、個人隱私 PII、未公開財務/併購、弱掃滲透細節、核心 IP\n• 建立第三方 AI 模型與供應商安全風險評估機制",
     COLOR_CYAN),
    ("🔒 2. 生成式 AI DLP 與安全 RAG 知識庫",
     "• 部署企業級 ChatGPT Enterprise / Codex 受控環境\n• 導入 GenAI 資料防洩漏 (DLP) 即時遮罩與 Token 稽核\n• 實施 Human-in-the-Loop 人機複核：敏感指令強制主管審批，AI 不具直接修改生產 DB 權限",
     COLOR_EMERALD),
    ("⚙️ 3. AI in SDLC (AISDLC) 研發雙軌治理",
     "• AI 輔助代碼審查 (Code Review) 與單元測試生成\n• 整合 SAST/DAST 品質閘門，開發交付週期縮短 40%\n• 建立 Prompt 風控、幻覺抑制與活動日誌 SIEM 串接，達到 100% 可觀測可調查",
     COLOR_PURPLE)
]

left_pos = 0.8
for title, desc, color in ai_pillars:
    box = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left_pos), Inches(1.6), Inches(3.75), Inches(5.3))
    box.fill.solid(); box.fill.fore_color.rgb = BG_CARD; box.line.color.rgb = color; box.line.width = Pt(1.5)
    
    tb = slide6.shapes.add_textbox(Inches(left_pos + 0.15), Inches(1.8), Inches(3.45), Inches(4.9))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = title; p.font.name = FONT_TITLE; p.font.size = Pt(13.5); p.font.bold = True; p.font.color.rgb = color
    
    p_d = tf.add_paragraph()
    p_d.text = desc; p_d.font.name = FONT_BODY; p_d.font.size = Pt(11); p_d.font.color.rgb = TEXT_LIGHT; p_d.space_before = Pt(10)
    left_pos += 3.98

# ==========================================
# SLIDE 7: 成果與商業效益 (ROI & 權威驗證)
# ==========================================
slide7 = prs.slides.add_slide(slide_layout)
set_slide_background(slide7)
add_header(slide7, "六、最終成果與商業價值：卓越的量化效益與權威實名佐證", "PROVEN BUSINESS OUTCOMES & RECOGNITION")

kpis = [
    ("100%", "Zero Outage 零停機紀錄", "GKE 多雲高可用跨區容災，春節促銷流量暴增下 100% 穩定連線", COLOR_CYAN),
    ("-30%", "多雲 FinOps 成本壓降", "精準實施資源彈性調度與節點治理，每年節省數百萬美元 TCO", COLOR_EMERALD),
    ("-30%", "重大資安事件大幅降低", "集中式 SOC/SIEM、EDR 聯防與零信任防護，杜絕重大勒索與資料外洩", COLOR_PURPLE),
    ("-30%", "平均修復時間 (MTTR)", "全域可觀測性與即時關聯分析，故障通報與排查縮短至分鐘級", COLOR_GOLD)
]

left_pos = 0.8
for val, title, desc, color in kpis:
    box = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left_pos), Inches(1.6), Inches(2.78), Inches(2.3))
    box.fill.solid(); box.fill.fore_color.rgb = BG_CARD; box.line.color.rgb = color; box.line.width = Pt(1.5)
    
    tb = slide7.shapes.add_textbox(Inches(left_pos + 0.1), Inches(1.7), Inches(2.58), Inches(2.1))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = val; p.font.name = FONT_TITLE; p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = color; p.alignment = PP_ALIGN.CENTER
    
    p_t = tf.add_paragraph()
    p_t.text = title; p_t.font.name = FONT_TITLE; p_t.font.size = Pt(11); p_t.font.bold = True; p_t.font.color.rgb = TEXT_WHITE; p_t.alignment = PP_ALIGN.CENTER
    
    p_d = tf.add_paragraph()
    p_d.text = desc; p_d.font.name = FONT_BODY; p_d.font.size = Pt(9.5); p_d.font.color.rgb = TEXT_MUTED; p_d.alignment = PP_ALIGN.CENTER; p_d.space_before = Pt(3)
    left_pos += 2.97

rec_box = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.1), Inches(11.7), Inches(2.8))
rec_box.fill.solid(); rec_box.fill.fore_color.rgb = BG_CARD_LIGHT; rec_box.line.color.rgb = COLOR_CYAN

tb_rec = slide7.shapes.add_textbox(Inches(1.0), Inches(4.2), Inches(11.3), Inches(2.6))
tf_rec = tb_rec.text_frame; tf_rec.word_wrap = True

p = tf_rec.paragraphs[0]; p.text = "🏆 業界高度認可與官方權威驗證出處 (Industry Citations)"; p.font.name = FONT_TITLE; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = COLOR_GOLD

recs = [
    "🎖️ 雲端架構卓越獎 (Cloud Architecture Excellence Award)：表彰高可用性、營運韌性與 FinOps 成本最佳化之典範實踐。",
    "📹 Google Cloud 官方全球客戶成功案例：親自接受 Google Cloud APAC 影音專訪，公開剖析 GKE 零停機架構心得。",
    "📰 主流科技媒體封面報導：CIO Taiwan 總編輯專題專訪《善用公有雲服務 搶攻全球商機》；iThome《運用 MongoDB Atlas 建構高可用平台》。",
    "🎤 產官學重量級 Keynote：受邀擔任 MongoDB.local Taipei 主講、CIO 價值學院 Keynote、行政院人事行政總處大師講座講師。"
]
for r in recs:
    p = tf_rec.add_paragraph()
    p.text = r
    p.font.name = FONT_BODY; p.font.size = Pt(10.5); p.font.color.rgb = TEXT_LIGHT
    p.space_after = Pt(4)

# ==========================================
# SLIDE 8: 願景與對貴公司的價值賦能
# ==========================================
slide8 = prs.slides.add_slide(slide_layout)
set_slide_background(slide8)
add_header(slide8, "七、未來展望：為 貴公司 創造的三大核心價值承諾", "STRATEGIC VALUE PROPOSITION FOR YOUR COMPANY")

commitments = [
    ("🏛️ 1. 董事會級戰略對齊與 ROI 驅動", 
     "• 絕不盲目堆疊技術，以 Business 策略為依歸\n• 建立可量化的資安與 IT 投資回報矩陣 (ROI/TCO)\n• 將資安態勢轉化為清晰的高階決策儀表板，助攻業務全球擴張",
     COLOR_CYAN),
    ("🛡️ 2. 構建全方位數位信任與極致韌性", 
     "• 落地 Zero Trust 零信任架構與抗勒索防禦縱深\n• 確保核心業務系統 99.99%+ 高可用與 100% 稽核即時合規 (Audit-Ready)\n• 守護核心機密紀錄、客戶隱私與智慧財產權",
     COLOR_EMERALD),
    ("🤖 3. 引領 AI 安全賦能與營運數位轉型", 
     "• 導入 ISO/IEC 42001 健全框架，安全普及生成式 AI 應用\n• 推動 DevSecOps 與 AISDLC 研發提效 35%+\n• 打造敏捷、高效且具備高度當責精神的現代化科技團隊",
     COLOR_PURPLE)
]

left_pos = 0.8
for title, desc, color in commitments:
    box = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left_pos), Inches(1.6), Inches(3.75), Inches(5.3))
    box.fill.solid(); box.fill.fore_color.rgb = BG_CARD; box.line.color.rgb = color; box.line.width = Pt(1.5)
    
    tb = slide8.shapes.add_textbox(Inches(left_pos + 0.15), Inches(1.8), Inches(3.45), Inches(4.9))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = title; p.font.name = FONT_TITLE; p.font.size = Pt(13.5); p.font.bold = True; p.font.color.rgb = color
    
    p_d = tf.add_paragraph()
    p_d.text = desc; p_d.font.name = FONT_BODY; p_d.font.size = Pt(11); p_d.font.color.rgb = TEXT_LIGHT; p_d.space_before = Pt(10)
    left_pos += 3.98

# ==========================================
# SLIDE 9: 致謝與 Q&A 交流
# ==========================================
slide9 = prs.slides.add_slide(slide_layout)
set_slide_background(slide9)

center_box = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.0), Inches(1.2), Inches(9.333), Inches(5.1))
center_box.fill.solid(); center_box.fill.fore_color.rgb = BG_CARD; center_box.line.color.rgb = BORDER_CYAN; center_box.line.width = Pt(1.5)

tb_end = slide9.shapes.add_textbox(Inches(2.5), Inches(1.6), Inches(8.333), Inches(4.3))
tf_end = tb_end.text_frame; tf_end.word_wrap = True

p1 = tf_end.paragraphs[0]
p1.text = "感謝各位 董事長、總經理／CEO 及高階長官"
p1.font.name = FONT_TITLE; p1.font.size = Pt(26); p1.font.bold = True; p1.font.color.rgb = TEXT_WHITE; p1.alignment = PP_ALIGN.CENTER

p2 = tf_end.add_paragraph()
p2.text = "Q & A  交流與討論"
p2.font.name = FONT_TITLE; p2.font.size = Pt(20); p2.font.bold = True; p2.font.color.rgb = COLOR_GOLD; p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(8)

p3 = tf_end.add_paragraph()
p3.text = "「科技引領創新，資安捍衛信任，策略創造價值。」"
p3.font.name = FONT_BODY; p3.font.size = Pt(14); p3.font.italic = True; p3.font.color.rgb = COLOR_CYAN; p3.alignment = PP_ALIGN.CENTER; p3.space_before = Pt(16)

p4 = tf_end.add_paragraph()
p4.text = "廖倫豪 博士 (Howard Liao, Ph.D.)\n集團資安長 暨 科技副總 (Group CISO)\n📱 +886-975-323161  |  ✉️ Liao.Howard@gmail.com\n🌐 https://howardliao.github.io/portfolio/  |  🔗 https://howardliao.github.io/Howard_CISO/"
p4.font.name = FONT_BODY; p4.font.size = Pt(11); p4.font.color.rgb = TEXT_LIGHT; p4.alignment = PP_ALIGN.CENTER; p4.space_before = Pt(20)

out_pptx_ciso = os.path.join(target_dir, "Howard_Liao_CISO_10Min_Executive_Presentation.pptx")
out_pptx_root = "/Users/howardliao/Howard_Liao_CISO_10Min_Executive_Presentation.pptx"

prs.save(out_pptx_ciso)
prs.save(out_pptx_root)

print(f"Successfully generated 10-minute executive presentation at:\n1. {out_pptx_ciso}\n2. {out_pptx_root}")
